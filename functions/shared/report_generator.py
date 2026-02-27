"""
Executive report generator — PDF and PPTX formats.

Generates CISO/CIO-ready reports from aggregated posture data
and AI-generated insights.
"""

import io
import json
import logging
from datetime import datetime, timezone

from shared.ai_agent import ask_executive_agent
from shared.normalizer import compute_statewide_aggregates
from shared.table_client import read_assessment_summaries, read_latest_snapshots_all_agencies

log = logging.getLogger(__name__)


def _get_report_data(agency_filter: str | None = None) -> dict:
    """Gather all data needed for the executive report."""
    snapshots = read_latest_snapshots_all_agencies()
    if agency_filter:
        snapshots = [s for s in snapshots if s.get("PartitionKey") == agency_filter]

    aggregates = compute_statewide_aggregates(snapshots)
    assessments = read_assessment_summaries(agency_filter)

    # Get AI-generated executive summary
    ai_result = ask_executive_agent(
        "Generate an executive summary of the current statewide compliance posture. "
        "Include key findings, highest-concern agencies, and prioritized recommendations.",
        agency_filter,
    )

    return {
        "generated_at": datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC"),
        "aggregates": aggregates,
        "snapshots": sorted(snapshots, key=lambda s: s.get("ComplianceScorePct", 0)),
        "assessments": assessments,
        "ai_summary": ai_result.get("answer", ""),
    }


def generate_pdf(agency_filter: str | None = None) -> bytes:
    """Generate a PDF executive summary report.

    Returns:
        PDF file content as bytes.
    """
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    data = _get_report_data(agency_filter)
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, topMargin=0.75 * inch)
    styles = getSampleStyleSheet()
    story = []

    # Title
    title_style = ParagraphStyle("Title", parent=styles["Title"], fontSize=18, spaceAfter=12)
    story.append(Paragraph("Purview Governance — Executive Summary", title_style))
    story.append(Paragraph(f"Generated: {data['generated_at']}", styles["Normal"]))
    story.append(Spacer(1, 0.3 * inch))

    # Statewide metrics
    agg = data.get("aggregates", {})
    if agg:
        story.append(Paragraph("Statewide Compliance Posture", styles["Heading2"]))
        metrics_data = [
            ["Metric", "Value"],
            ["Total Agencies", str(agg.get("TotalAgencies", 0))],
            ["Avg Compliance Score", f"{agg.get('AvgComplianceScore', 0):.1f}%"],
            ["Min / Max Score", f"{agg.get('MinComplianceScore', 0):.1f}% / {agg.get('MaxComplianceScore', 0):.1f}%"],
            ["Avg Label Coverage", f"{agg.get('AvgLabelCoverage', 0):.1f}%"],
            ["DLP Incidents (30d)", str(agg.get("TotalDlpIncidents30d", 0))],
            ["Ext. Sharing Events", str(agg.get("TotalExternalSharing", 0))],
            ["Insider Risk Alerts", str(agg.get("TotalInsiderRiskAlerts", 0))],
        ]
        t = Table(metrics_data, colWidths=[3 * inch, 3 * inch])
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1a3c6e")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
            ("ALIGN", (1, 0), (1, -1), "CENTER"),
        ]))
        story.append(t)
        story.append(Spacer(1, 0.3 * inch))

    # AI Summary
    if data.get("ai_summary"):
        story.append(Paragraph("AI Executive Analysis", styles["Heading2"]))
        for para in data["ai_summary"].split("\n\n"):
            if para.strip():
                story.append(Paragraph(para.strip(), styles["Normal"]))
                story.append(Spacer(1, 6))
        story.append(Spacer(1, 0.2 * inch))

    # Per-agency table
    snapshots = data.get("snapshots", [])
    if snapshots:
        story.append(Paragraph("Agency Detail", styles["Heading2"]))
        agency_data = [["Agency", "Compliance %", "Label Coverage %", "DLP (30d)", "Ext. Sharing"]]
        for s in snapshots[:20]:
            agency_data.append([
                s.get("PartitionKey", ""),
                f"{s.get('ComplianceScorePct', 0):.1f}%",
                f"{s.get('LabelCoveragePct', 0):.1f}%",
                str(s.get("DlpIncidents30d", 0)),
                str(s.get("ExternalSharingCount", 0)),
            ])
        t = Table(agency_data)
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1a3c6e")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
            ("ALIGN", (1, 0), (-1, -1), "CENTER"),
        ]))
        story.append(t)

    # Footer
    story.append(Spacer(1, 0.5 * inch))
    footer_style = ParagraphStyle("Footer", parent=styles["Normal"], fontSize=7, textColor=colors.grey)
    story.append(Paragraph(
        "This report contains metadata only — no document content, PII, or user identities. "
        "All scores are native from Microsoft Purview and Compliance Manager.",
        footer_style,
    ))

    doc.build(story)
    return buffer.getvalue()


def generate_pptx(agency_filter: str | None = None) -> bytes:
    """Generate a PowerPoint executive summary.

    Returns:
        PPTX file content as bytes.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    data = _get_report_data(agency_filter)
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Purview Governance\nExecutive Summary"
    slide.placeholders[1].text = f"Generated: {data['generated_at']}\nMetadata Only — No PII"

    # Statewide metrics slide
    agg = data.get("aggregates", {})
    if agg:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Statewide Compliance Posture"
        body = slide.placeholders[1].text_frame
        body.text = f"Agencies Reporting: {agg.get('TotalAgencies', 0)}"
        body.add_paragraph().text = f"Average Compliance Score: {agg.get('AvgComplianceScore', 0):.1f}%"
        body.add_paragraph().text = f"Score Range: {agg.get('MinComplianceScore', 0):.1f}% – {agg.get('MaxComplianceScore', 0):.1f}%"
        body.add_paragraph().text = f"Average Label Coverage: {agg.get('AvgLabelCoverage', 0):.1f}%"
        body.add_paragraph().text = f"DLP Incidents (30d): {agg.get('TotalDlpIncidents30d', 0)}"
        body.add_paragraph().text = f"Lowest Score Agency: {agg.get('LowestComplianceAgency', 'N/A')}"

    # AI Summary slide
    if data.get("ai_summary"):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "AI Executive Analysis"
        body = slide.placeholders[1].text_frame
        # Split into paragraphs, limit to fit slide
        paragraphs = [p.strip() for p in data["ai_summary"].split("\n") if p.strip()]
        body.text = paragraphs[0] if paragraphs else ""
        for para in paragraphs[1:15]:  # Limit to prevent overflow
            body.add_paragraph().text = para

    # Footer slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = (
        "This report contains metadata only — no document content, PII, or user identities.\n"
        "All scores are native from Microsoft Purview and Compliance Manager."
    )
    tf.paragraphs[0].font.size = Pt(10)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
